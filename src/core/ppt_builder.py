"""PPT building module for creating PowerPoint presentations."""

import io
from pathlib import Path
from typing import List, Tuple, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Length, Pt

from .pdf_parser import FontInfo


# Conversion ratio: 1 inch = 914400 EMU, 1 inch = 72 points
# PDF uses points, PPT uses EMU (English Metric Units)
PDF_TO_EMU = 12700  # 914400 / 72


class PPTBuilder:
    """Build PowerPoint presentations from extracted PDF elements."""

    def __init__(self, pdf_width: float, pdf_height: float):
        """
        Initialize the PPT builder.

        Args:
            pdf_width: PDF page width in points
            pdf_height: PDF page height in points
        """
        self.prs = Presentation()
        self.pdf_width = pdf_width
        self.pdf_height = pdf_height

        # Set slide size to match PDF page size
        self._set_slide_size()

    def _set_slide_size(self):
        """Set the slide size to match the PDF page dimensions."""
        # Convert PDF points to inches (72 points = 1 inch)
        width_inches = self.pdf_width / 72
        height_inches = self.pdf_height / 72

        self.prs.slide_width = Inches(width_inches)
        self.prs.slide_height = Inches(height_inches)

    def add_slide(self):
        """Add a new blank slide to the presentation."""
        # Use blank layout
        blank_layout = self.prs.slide_layouts[6]
        self.prs.slides.add_slide(blank_layout)

    def add_text_box(self, text: str, bbox: Tuple[float, float, float, float], font_info: FontInfo):
        """
        Add a text box to the current slide.

        Args:
            text: Text content
            bbox: Bounding box (x0, y0, x1, y1) in PDF points
            font_info: Font information
        """
        if not self.prs.slides:
            self.add_slide()

        slide = self.prs.slides[-1]

        # Calculate dimensions in inches
        left, top, width, height = self._bbox_to_inches(bbox)

        # Add text box
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Set text content
        p = text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.LEFT

        # Apply font formatting
        font = p.font
        font.name = font_info.font or "Calibri"
        font.size = Pt(font_info.size)

        # Set color - detect if color is too light for white background
        if isinstance(font_info.color, (tuple, list)) and len(font_info.color) >= 3:
            r, g, b = font_info.color[:3]
            # Calculate luminance to detect if text is too light for white background
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            # If luminance > 0.5, text is too light for white background, use dark gray
            if luminance > 0.5:
                r, g, b = 51, 51, 51  # Dark gray for better readability
            font.color.rgb = RGBColor(r, g, b)
        else:
            # Default to black if no color specified
            font.color.rgb = RGBColor(0, 0, 0)

        # Set bold/italic
        font.bold = font_info.bold
        font.italic = font_info.italic

        # Auto-size the text frame to fit content
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    def add_image(self, image_bytes: bytes, bbox: Tuple[float, float, float, float]):
        """
        Add an image to the current slide.

        Args:
            image_bytes: Image data as bytes
            bbox: Bounding box (x0, y0, x1, y1) in PDF points
        """
        if not self.prs.slides:
            self.add_slide()

        slide = self.prs.slides[-1]

        # Calculate dimensions in inches
        left, top, width, height = self._bbox_to_inches(bbox)

        # Create image stream
        image_stream = io.BytesIO(image_bytes)

        # Add image to slide
        slide.shapes.add_picture(image_stream, Inches(left), Inches(top), Inches(width), Inches(height))

    def add_table(
        self, rows: int, cols: int, bbox: Tuple[float, float, float, float], data: List[List[str]]
    ):
        """
        Add a table to the current slide.

        Args:
            rows: Number of rows
            cols: Number of columns
            bbox: Bounding box (x0, y0, x1, y1) in PDF points
            data: Table data as 2D list
        """
        if not self.prs.slides or rows == 0 or cols == 0:
            return

        slide = self.prs.slides[-1]

        # Calculate dimensions in inches
        left, top, width, height = self._bbox_to_inches(bbox)

        # Calculate cell dimensions
        cell_width = width / cols
        cell_height = height / rows

        # Add table
        table = slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        ).table

        # Fill table with data
        for i, row_data in enumerate(data):
            if i >= rows:
                break
            for j, cell_text in enumerate(row_data):
                if j >= cols:
                    break

                cell = table.rows[i].cells[j]
                cell.text = str(cell_text)

                # Set font size and alignment
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.alignment = PP_ALIGN.LEFT

    def _bbox_to_inches(self, bbox: Tuple[float, float, float, float]) -> Tuple[float, float, float, float]:
        """
        Convert PDF bounding box (points) to PPT dimensions (inches).

        Args:
            bbox: Bounding box (x0, y0, x1, y1) in PDF points

        Returns:
            Tuple of (left, top, width, height) in inches
        """
        x0, y0, x1, y1 = bbox

        # Convert to inches (72 points = 1 inch)
        left = x0 / 72
        top = y0 / 72
        width = (x1 - x0) / 72
        height = (y1 - y0) / 72

        return (left, top, width, height)

    def save(self, output_path: Path | str):
        """
        Save the presentation to a file.

        Args:
            output_path: Path to save the PPTX file
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)

    def get_bytes(self) -> bytes:
        """
        Get the presentation as bytes.

        Returns:
            Presentation data as bytes
        """
        stream = io.BytesIO()
        self.prs.save(stream)
        return stream.getvalue()
